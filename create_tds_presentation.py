from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_tds_presentation():
    prs = Presentation()
    
    # Set slide dimensions (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "TDS Compliance Implementation\nfor Loyalty Program"
    subtitle.text = "Executive Presentation for Senior Management\nEnsuring Regulatory Compliance under ITR Rules"
    
    # Slide 2: Executive Summary
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Executive Summary"
    content.text = ("The Challenge:\n"
                   "• ITR rules mandate TDS on loyalty rewards > ₹20,000 annually\n"
                   "• Current status: Zero TDS compliance - significant regulatory risk\n"
                   "• Solution: Automated TDS withholding system\n\n"
                   "Key Solution Benefits:\n"
                   "✓ Full regulatory compliance\n"
                   "✓ Automated TDS calculation & withholding\n"
                   "✓ Customer transparency\n"
                   "✓ Seamless redemption experience")
    
    # Slide 3: Regulatory Context
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Regulatory Context - Section 194R"
    content.text = ("Income Tax Rules Requirements:\n\n"
                   "• Threshold: ₹20,000 per financial year\n"
                   "• TDS Rate: 10% (PAN verified) / 20% (Non-PAN)\n"
                   "• Effective: Already in force\n\n"
                   "Penalty for Non-Compliance:\n"
                   "• Interest @ 1.5% per month\n"
                   "• Penalty up to TDS amount\n"
                   "• Prosecution proceedings possible")
    
    # Slide 4: Solution Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Our Solution: Three-Tier Limit System"
    content.text = ("PAN Verified Customers:\n"
                   "• Limit 1: ₹0 - No TDS\n"
                   "• Limit 2: ₹18,000 - Buffer zone monitoring\n"
                   "• Limit 3: ₹20,000 - TDS trigger point\n\n"
                   "Non-PAN Verified Customers:\n"
                   "• Limit 1: ₹0 - No benefits\n"
                   "• Limit 2: ₹16,000 - Earlier intervention\n"
                   "• Limit 3: ₹20,000 - Higher TDS rate (20%)")
    
    # Slide 5: Live Example
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Solution Logic - Live Example"
    content.text = ("Customer Approaching Threshold:\n\n"
                   "Initial State:\n"
                   "• Cashback Received (FY): ₹19,000\n"
                   "• Unredeemed Points: ₹500\n"
                   "• Customer earns: ₹1,800\n\n"
                   "Calculation:\n"
                   "• Total exposure: ₹21,300\n"
                   "• Exceeds ₹20,000 by: ₹1,300\n"
                   "• TDS @ 10%: ₹2,130\n"
                   "• TDS to deduct now: ₹1,630\n"
                   "• Redeemable points: ₹170")
    
    # Slide 6: Customer Journey
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Customer Journey"
    content.text = ("Before ₹20,000 Threshold:\n"
                   "1. Customer earns rewards normally\n"
                   "2. Full redemption available\n"
                   "3. System monitors cumulative total\n"
                   "4. Alert at ₹18,000 (PAN verified)\n\n"
                   "At/After ₹20,000 Threshold:\n"
                   "1. System calculates TDS automatically\n"
                   "2. Withholds TDS from current transaction\n"
                   "3. Shows net redeemable amount\n"
                   "4. Issues TDS certificate quarterly")
    
    # Slide 7: Implementation Roadmap
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Implementation Roadmap"
    content.text = ("Phase 1: Backend Development (Week 1-2)\n"
                   "• TDS calculation engine\n"
                   "• Database schema updates\n"
                   "• PAN verification integration\n\n"
                   "Phase 2: Frontend Integration (Week 3)\n"
                   "• Customer notification system\n"
                   "• Balance display updates\n\n"
                   "Phase 3: Compliance Setup (Week 4)\n"
                   "• Government portal integration\n"
                   "• Form 26AS generation\n\n"
                   "Phase 4: Launch (Month 2)\n"
                   "• Pilot with select customers\n"
                   "• Full rollout")
    
    # Slide 8: Financial Impact
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Financial Impact Analysis"
    content.text = ("Cost of Implementation:\n"
                   "• Development: ₹5,00,000\n"
                   "• Testing & QA: ₹1,00,000\n"
                   "• Integration: ₹2,00,000\n"
                   "• Total One-time: ₹8,00,000\n"
                   "• Recurring: ₹6,00,000/year\n\n"
                   "Cost of Non-Compliance:\n"
                   "• Penalty (1 year): ₹50,00,000+\n"
                   "• Interest: ₹9,00,000+\n"
                   "• Legal proceedings: ₹10,00,000+\n"
                   "• Total Risk: ₹69,00,000+\n\n"
                   "ROI: Cost recovered by avoiding just 2% of penalty risk")
    
    # Slide 9: Success Metrics
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Success Metrics & KPIs"
    content.text = ("Compliance Metrics:\n"
                   "• 100% TDS collection on applicable rewards\n"
                   "• 0% penalty from tax authorities\n"
                   "• 100% timely government remittance\n\n"
                   "Customer Metrics:\n"
                   "• <5% increase in support queries\n"
                   "• >90% customer satisfaction maintained\n"
                   "• <2% churn in high-value segment\n\n"
                   "Operational Metrics:\n"
                   "• 100% automated calculations\n"
                   "• <1% manual interventions\n"
                   "• 99.9% system uptime")
    
    # Slide 10: Risk Mitigation
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Risk Mitigation Strategy"
    content.text = ("Technical Risks:\n"
                   "• Calculation errors → Extensive testing, audit trails\n"
                   "• System downtime → Redundancy, offline capability\n"
                   "• Integration failure → Fallback mechanisms\n\n"
                   "Customer Risks:\n"
                   "• Confusion about TDS → Clear communication\n"
                   "• Redemption issues → Grace period, manual override\n"
                   "• Certificate delays → Automated generation\n\n"
                   "Compliance Risks:\n"
                   "• Regular audits and reconciliation\n"
                   "• Legal team oversight\n"
                   "• Government liaison established")
    
    # Slide 11: Competitive Advantage
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Competitive Advantage"
    content.text = ("Industry Comparison:\n"
                   "• Competitor A: Manual TDS, customer complaints\n"
                   "• Competitor B: No TDS, under investigation\n"
                   "• Our Solution: Automated, compliant, transparent\n\n"
                   "First-Mover Benefits:\n"
                   "✓ Regulatory goodwill\n"
                   "✓ Customer trust\n"
                   "✓ Industry benchmark\n"
                   "✓ Potential to white-label solution\n\n"
                   "Long-term Value:\n"
                   "• Foundation for comprehensive tax management\n"
                   "• Scalable to other reward programs")
    
    # Slide 12: Recommendations
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Recommendations & Next Steps"
    content.text = ("Immediate Actions Required:\n"
                   "1. Approve project budget: ₹8,00,000\n"
                   "2. Form implementation team\n"
                   "3. Set target date: Before March 31st\n"
                   "4. Initiate customer communication\n\n"
                   "Decision Required:\n"
                   "Option A: Full Implementation (Recommended)\n"
                   "• Complete automation, 6 weeks, ₹8,00,000\n\n"
                   "Option B: Phased Approach\n"
                   "• High-value customers first, 3 months, ₹5,00,000\n\n"
                   "Recommendation: Option A - Compliance risk too high to delay")
    
    # Slide 13: Technical Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Technical Architecture"
    content.text = ("System Flow:\n"
                   "1. Real-time Monitoring - Track cumulative rewards\n"
                   "2. Threshold Detection - Automatic triggers\n"
                   "3. TDS Calculation - Based on PAN status\n"
                   "4. Withholding - Deduct from transaction\n"
                   "5. Remittance - Monthly to government\n"
                   "6. Certification - Quarterly to customers\n\n"
                   "Database Changes:\n"
                   "• New fields: TDS_held, TDS_applicable, PAN_verified\n"
                   "• New tables: TDS_transactions, TDS_certificates\n"
                   "• Complete audit trail for compliance")
    
    # Slide 14: Q&A
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Questions & Discussion"
    content.text = ("Key Discussion Points:\n\n"
                   "• Timeline flexibility?\n"
                   "• Budget approval process?\n"
                   "• Legal team involvement?\n"
                   "• Customer segment prioritization?\n"
                   "• Communication strategy approval?\n\n"
                   "Contact: Project Management Team\n"
                   "Next Review: [To be scheduled]")
    
    # Save presentation
    prs.save('TDS_Compliance_Presentation.pptx')
    print("PowerPoint presentation created successfully!")
    print("File saved as: TDS_Compliance_Presentation.pptx")

# Run the function
create_tds_presentation()